import os
import re
import fitz
import io
import pytesseract
from PIL import Image
from pptx import Presentation
from docx import Document
import openpyxl
from langchain_text_splitters import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct

# tesseract path for windows
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

collection_name = "pdf_ppt_xl"
vector_size     = 384
batch_size      = 32
chunk_size      = 1000
chunk_overlap   = 200


# symbol removal and text cleaning
def clean_text(text):
    text = re.sub(r'\x00', ' ', text)
    text = re.sub(r'[\u2500-\u27FF]', ' ', text)
    text = re.sub(r'[\u2000-\u206F]', ' ', text)
    text = re.sub(r'[^\x20-\x7E]', ' ', text)
    text = re.sub(r'\n+', ' ', text)
    text = re.sub(r' +', ' ', text)
    return text.strip()


# garbage detection after cleaning
def is_garbage_text(text):
    cleaned = re.sub(r'[^\x20-\x7E]', ' ', text)
    cleaned = re.sub(r' +', ' ', cleaned).strip()
    words   = cleaned.split()

    if len(words) < 10:
        return True

    avg_len = sum(len(w) for w in words) / len(words)
    if avg_len < 2.5:
        return True

    return False


# load pdf file page by page using fitz
def load_pdf(file_path):
    filename = os.path.basename(file_path)
    doc      = fitz.open(file_path)
    pages    = []
    skipped  = 0

    print(f"Opening PDF '{filename}' {len(doc)} pages")

    for page_num in range(len(doc)):
        raw_text = doc[page_num].get_text()

        if is_garbage_text(raw_text):
            skipped += 1
            continue

        cleaned = clean_text(raw_text)

        if len(cleaned) < 50:
            skipped += 1
            continue

        pages.append({
            "text"    : cleaned,
            "page_num": page_num + 1,
            "source"  : filename
        })

    doc.close()
    print(f"  Loaded {len(pages)} pages | Skipped {skipped}")
    return pages


# load image file using pytesseract ocr
def load_image(file_path):
    filename = os.path.basename(file_path)
    print(f"Opening Image '{filename}'")

    image = Image.open(file_path)
    text  = pytesseract.image_to_string(image, lang="eng")

    if is_garbage_text(text):
        print(f"  No readable text found in image")
        return []

    cleaned = clean_text(text)

    if len(cleaned) < 50:
        return []

    # image has no pages so page_num is 1
    pages = [{
        "text"    : cleaned,
        "page_num": 1,
        "source"  : filename
    }]

    print(f"  Loaded 1 page from image")
    return pages

def load_docx(file_path):
    filename = os.path.basename(file_path)
    doc      = Document(file_path)
    pages    = []
    para_batch = []
    page_num   = 1

    print(f"Opening DOCX '{filename}'")

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            para_batch.append(text)

        # every 20 paragraphs treat as one page
        if len(para_batch) >= 20:
            full_text = " ".join(para_batch)
            cleaned   = clean_text(full_text)

            if not is_garbage_text(full_text) and len(cleaned) >= 50:
                pages.append({
                    "text"    : cleaned,
                    "page_num": page_num,
                    "source"  : filename
                })
                page_num += 1

            para_batch = []

    # remaining paragraphs
    if para_batch:
        full_text = " ".join(para_batch)
        cleaned   = clean_text(full_text)

        if not is_garbage_text(full_text) and len(cleaned) >= 50:
            pages.append({
                "text"    : cleaned,
                "page_num": page_num,
                "source"  : filename
            })

    print(f"  Loaded {len(pages)} sections from docx")
    return pages

# load ppt file slide by slide using python-pptx
def extract_images_text_from_slide(slide):
    """
    extracts all images from a ppt slide
    runs ocr on each image
    returns combined text from all images
    """
    image_texts = []
    
    for shape in slide.shapes:
        # check if shape is a picture
        if shape.shape_type != 13:
            continue
            
        try:
            # get image bytes from shape
            img_bytes = shape.image.blob
            image     = Image.open(io.BytesIO(img_bytes))
            ocr_text  = pytesseract.image_to_string(image, lang="eng")
            
            if ocr_text.strip():
                image_texts.append(ocr_text.strip())
                
        except Exception:
            continue
    
    return " ".join(image_texts)


def load_ppt(file_path):
    filename = os.path.basename(file_path)
    prs      = Presentation(file_path)
    pages    = []
    skipped  = 0

    print(f"Opening PPT '{filename}' {len(prs.slides)} slides")

    for slide_num, slide in enumerate(prs.slides):
        slide_text = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                # join runs and strip
                line = " ".join(run.text for run in para.runs).strip()
                # clean multiple spaces inside line
                line = " ".join(line.split())
                if line:
                    slide_text.append(line)

        # get image text from slide
        image_text = extract_images_text_from_slide(slide)

        # combine text and image text
        full_text  = " ".join(slide_text) + " " + image_text
        full_text  = " ".join(full_text.split())  # clean all extra spaces

        # skip only if truly empty after combining both
        if len(full_text.strip()) < 3:
            skipped += 1
            continue

        cleaned = clean_text(full_text)

        if len(cleaned.strip()) < 3:
            skipped += 1
            continue

        pages.append({
            "text"    : cleaned,
            "page_num": slide_num + 1,
            "source"  : filename
        })

    print(f"  Loaded {len(pages)} slides | Skipped {skipped}")
    return pages

# load xlsx file row by row using openpyxl


def load_xlsx(file_path):
    filename = os.path.basename(file_path)
    wb       = openpyxl.load_workbook(file_path, data_only=True)
    pages    = []

    print(f"Opening XLSX '{filename}' {len(wb.sheetnames)} sheets")

    for sheet_num, sheet_name in enumerate(wb.sheetnames):
        ws         = wb[sheet_name]
        sheet_rows = []

        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(
                str(cell) for cell in row
                if cell is not None and str(cell).strip()
            )
            if row_text.strip():
                sheet_rows.append(row_text)

        full_text = " ".join(sheet_rows)

        # only skip if truly empty
        if len(full_text.strip()) < 10:
            continue

        cleaned = clean_text(full_text)

        pages.append({
            "text"    : cleaned,
            "page_num": sheet_num + 1,
            "source"  : filename
        })

    print(f"  Loaded {len(pages)} sheets from xlsx")
    return pages


# detect file type and call the right loader
def load_file(file_path):
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return load_pdf(file_path)

    elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"]:
        return load_image(file_path)

    elif ext in [".pptx", ".ppt"]:
        return load_ppt(file_path)

    elif ext in [".xlsx", ".xls"]:
        return load_xlsx(file_path)

    elif ext in [".docx", ".doc"]:      # ← yeh add karo
        return load_docx(file_path)

    else:
        print(f"Unsupported file type: {ext}")
        return []

# recursive chunking same as before
def make_chunks(pages):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size     =chunk_size,
        chunk_overlap  =chunk_overlap,
        separators     =["\n\n", "\n", ". ", " ", ""],
        length_function=len
    )

    all_chunks = []
    chunk_id   = 0

    for page in pages:
        chunks = splitter.split_text(page["text"])

        for chunk in chunks:
            chunk = chunk.strip()

            if len(chunk) < 80:
                continue

            all_chunks.append({
                "chunk_text": chunk,
                "page_num"  : page["page_num"],
                "source"    : page["source"],
                "chunk_id"  : chunk_id
            })
            chunk_id += 1

    print(f"Total chunks: {len(all_chunks)}")
    return all_chunks


# model loading
def load_model():
    print("Loading model...")
    model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")
    print("Model ready")
    return model


# batch embeddings
def embed_chunks(chunks, model):
    texts = [chunk["chunk_text"] for chunk in chunks]
    print(f"Embedding {len(texts)} chunks")

    vectors = model.encode(
        texts,
        batch_size          =batch_size,
        show_progress_bar   =True,
        convert_to_numpy    =True,
        normalize_embeddings=True
    )

    print(f"Embeddings done shape {vectors.shape}")
    return vectors


# qdrant setup
def setup_qdrant():
    try:
        client = QdrantClient(host="localhost", port=6333)
        client.get_collections()
        print("Qdrant connected (Docker @ localhost:6333)")
        return client
    except Exception as e:
        print("Qdrant is not connected")
        print("First run: docker run -p 6333:6333 qdrant/qdrant")
        raise e


# create fresh collection every time
def create_collection(client):
    existing = [c.name for c in client.get_collections().collections]

    if collection_name in existing:
        client.delete_collection(collection_name=collection_name)
        print(f"Old collection '{collection_name}' deleted")

    client.create_collection(
        collection_name=collection_name,
        vectors_config =VectorParams(
            size    =vector_size,
            distance=Distance.COSINE
        )
    )
    print(f"Collection '{collection_name}' created")
    return True


# insert points into qdrant
def insert_to_qdrant(chunks, vectors, client):
    points = []

    for i, (chunk, vector) in enumerate(zip(chunks, vectors)):
        points.append(PointStruct(
            id     = i,
            vector = vector.tolist(),
            payload = {
                "chunk_text": chunk["chunk_text"],
                "page_num"  : chunk["page_num"],
                "source"    : chunk["source"],
                "chunk_id"  : chunk["chunk_id"]
            }
        ))

    client.upsert(
        collection_name=collection_name,
        points=points
    )
    print(f"Inserted {len(points)} points")


# verify insertion
def verify_insert(client):
    info   = client.get_collection(collection_name)
    sample = client.retrieve(
        collection_name=collection_name,
        ids            =[0],
        with_payload   =True,
        with_vectors   =False
    )

    print(f"\nVerification")
    print(f"  Points stored : {info.points_count}")
    print(f"  Vector size   : {info.config.params.vectors.size}")
    print(f"  chunk_text    : {sample[0].payload['chunk_text'][:200]}")
    print(f"  page_num      : {sample[0].payload['page_num']}")
    print(f"  source        : {sample[0].payload['source']}")


# main pipeline accepts single file or list of files
def build_knowledge_base(file_paths):
    client = setup_qdrant()
    create_collection(client)

    # accept single file or list of files
    if isinstance(file_paths, str):
        file_paths = [file_paths]

    all_pages  = []

    for file_path in file_paths:
        print(f"\nProcessing: {file_path}")
        pages = load_file(file_path)
        all_pages.extend(pages)

    print(f"\nTotal pages loaded: {len(all_pages)}")

    chunks  = make_chunks(all_pages)
    model   = load_model()
    vectors = embed_chunks(chunks, model)
    insert_to_qdrant(chunks, vectors, client)
    verify_insert(client)

    print(f"\nKnowledge base ready")
    print(f"  Total files   : {len(file_paths)}")
    print(f"  Total pages   : {len(all_pages)}")
    print(f"  Total chunks  : {len(chunks)}")
    print(f"  Vector shape  : {vectors.shape}")

    return client, model


if __name__ == "__main__":
    # single file
    # client, model = build_knowledge_base("GDG Inductions 2026.pdf")

    # multiple files
    files = [
        "GDG Inductions 2026.pdf",
        "Intern Essentials PClub PPT.pptx",
        "stu_marks.xlsx",
        "APP App Dev Workshop-3.png"
        # "notes.docx"
    ]
    client, model = build_knowledge_base(files)
